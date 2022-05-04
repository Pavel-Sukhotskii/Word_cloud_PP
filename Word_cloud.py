from wordcloud import WordCloud
from io import StringIO
import random
from deep_translator import GoogleTranslator

import matplotlib.pyplot as plt
import seaborn as sns

import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN

# nltk.download('stopwords') #IMPORTANT TO DOWNLOAD FIRST TIME
# nltk.download('punkt') #IMPORTANT TO DOWNLOAD FIRST TIME
stop_words = set(stopwords.words('english'))

IMGS_FOLDER_W = './imgs_words'

FILE_NAME = 'Word_clouds'

if not os.path.exists(IMGS_FOLDER_W):
    os.makedirs(IMGS_FOLDER_W)
    
# column_ids like: ['QOp0a', 'QOp0b']    

def word_cloud(logs, column_ids:list, translation=None, rows=False): # translation - language
    
    '''Creating pictures and titles'''
    
    save_folder = 'imgs_words/'
    
    titles = {}
    
    if rows == True:
        data = pd.DataFrame()
        # converting to type: name - column names, answers - rows
        for q_id in column_ids:
            name = logs.question[logs.question.str.contains(q_id)].iloc[0]
            data[name] = logs[logs.question == name].answer
    else:
        data = logs
            
    for id in column_ids:
        
        column_name = data.columns[data.columns.str.contains(id)][0]
        
        titles[id] = column_name.split(id+'. ')[1] # split to get only question
        
        words = data[column_name].dropna().str.lower() # getting column with answers
        if translation:
            words = words.apply(lambda x: GoogleTranslator(source=translation, target='english').translate(x)) # translation 
            
            # titles manual translation
            dict_to_translate ={
                'Puoi dirci cosa ti è piaciuto?':
                'Can you please let us know what you liked about it?',
                "C'è qualcosa che non hai apprezzato?": 'And is there anything you disliked about it?'
            }
            
            titles[id] = dict_to_translate[column_name.split(id+'. ')[1]]
            
            # titles automatical translation
#             titles[id] = GoogleTranslator(source=translation, target='english').translate(column_name.split(id+'. ')[1])
            
        text = ' '.join(list(words)).replace('-', '').replace('  ', '') # joining the answers to one str
        
        wordcloud = WordCloud().generate(text)
        wordcloud = WordCloud(width=1200, height=600, background_color="white",
                              max_words=len(text),max_font_size=210, relative_scaling=.01).generate(text) # WC params
        plt.clf()  
        plt.figure(figsize= [25, 15])
        plt.imshow(wordcloud)
        plt.axis("off")
        plt.savefig(save_folder + 'words_%s.png' % id, bbox_inches='tight', pad_inches=0)
        
        '''Creating histogram'''

        dict_to_hist = {}
        # creating a dict with counted words
        for word in [item for subset in [word_tokenize(i) for i in words] for item in subset]: # in list with tokenized words
            if len(word) < 3 or word in list(stop_words):
                continue
            if word in dict_to_hist.keys():
                dict_to_hist[word] += 1
            else:
                dict_to_hist[word] = 1

        dict_to_hist = dict(sorted(dict_to_hist.items(), key=lambda item: item[1])[::-1])   
        
        plt.clf()
        plt.figure(figsize= [25, 5])
        plt.grid(alpha=0.5)
        sns.barplot(x=list(dict_to_hist.keys())[:10], y=list(dict_to_hist.values())[:10])

        plt.savefig(save_folder + 'histogram_%s.png' % id, bbox_inches='tight', pad_inches=0)

    '''Creating presentation'''
    
    ppt = Presentation()

    ppt.slide_height=Cm(19.05)
    ppt.slide_width=Cm(33.867)
    
    for q_id, name in titles.items():
        
        slide_layout = ppt.slide_layouts[6] # 6 - empty list
        slide = ppt.slides.add_slide(slide_layout)
        
        '''Title'''
        width = Cm(15)
        height = Cm(1)
        left = Cm(9.43)
        top = Cm(0.3)

        title = slide.shapes.add_textbox(left, top, width, height)
        tf = title.text_frame
        tf.text = name
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER # align to center

        '''Word Cloud'''
        img_path = save_folder + 'words_%s.png' % q_id

        picture_height = Cm(11.664)
        picture_width = Cm(24.21)
        picture_top = Cm(1.69)
        picture_left = Cm(4.83)

        pic = slide.shapes.add_picture(img_path, picture_left, picture_top, height=picture_height, width=picture_width)
        
        '''Histogram'''
        hist_path = save_folder + 'histogram_%s.png' % q_id

        hist_height = Cm(4.93)
        hist_width = Cm(24.21)
        hist_top = Cm(13.74)
        hist_left = Cm(4.83)

        pic = slide.shapes.add_picture(hist_path, hist_left, hist_top, height=hist_height, width=hist_width)
        
    ppt.save('Results/%s.pptx' % FILE_NAME)        
